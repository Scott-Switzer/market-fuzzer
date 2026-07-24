"""Pitch deck builder: renders HTML + PPTX + PDF from REAL evidence.

Every number shown is read from artifacts/submission/<sha>/pitch/deck_data.json.
If no run exists, it fails loudly rather than fabricating data.

Guarantees enforced here (and audited by tests/submission_audit/test_deck_evidence.py):
  * exactly 10 slides in every format;
  * every metric comes from deck_data.json (no hand-entered figures);
  * every performance slide carries a TIER watermark derived from evidence metadata;
  * the deck only renders evidence for the CURRENT git SHA (no stale SHA);
  * only confirmed failures are shown (predicate-violating, still_fails-verified);
  * synthetic-fixture runs are never labeled "historical".
"""

from __future__ import annotations

import glob
import json
import subprocess
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

DECK_PATH = Path("app/static/pitch-deck/index.html")
DECK_DIR = Path("app/static/pitch-deck")

TAGLINE = "Backtesting shows where a strategy worked. Fenrix searches for where it breaks."
FLOW = "Describe \u2192 Review \u2192 Lock \u2192 Backtest \u2192 Stress \u2192 Minimize \u2192 Export"

TIER_LABELS = {
    1: "TIER 1 \u00b7 Fenrix anonymized bundle",
    2: "TIER 2 \u00b7 yfinance historical (actual run)",
    3: "TIER 3 \u00b7 deterministic synthetic fixture (NOT historical)",
}

HISTORICAL_MODES = {"yfinance", "fenrix"}


def _git_sha() -> str:
    try:
        out = subprocess.run(["git", "rev-parse", "HEAD"], capture_output=True, text=True, cwd=Path.cwd())
        return out.stdout.strip()[:16]
    except Exception:
        return "unknown"


def _latest_deck_data() -> dict:
    matches = sorted(glob.glob("artifacts/submission/*/pitch/deck_data.json"))
    if not matches:
        raise RuntimeError("no submission evidence found; run `make submission-demo` first")
    return json.loads(Path(matches[-1]).read_text())


# ---------------------------------------------------------------------------
# Evidence loading (current-SHA enforced)
# ---------------------------------------------------------------------------


@dataclass
class Evidence:
    data: dict
    base_dir: Path
    sha: str
    tier: int
    data_mode: str
    watermark: str
    screenshots: list[Path] = field(default_factory=list)
    equity_curve: list[float] = field(default_factory=list)


def load_evidence(require_current_sha: bool = True) -> Evidence:
    sha = _git_sha()
    base = Path(f"artifacts/submission/{sha}")
    deck_json = base / "pitch" / "deck_data.json"
    if not deck_json.exists():
        if require_current_sha:
            raise RuntimeError(
                f"no evidence for current git SHA {sha}; refusing stale deck. "
                "Run `make submission-demo` first."
            )
        data = _latest_deck_data()
        base = Path(f"artifacts/submission/{data['git_sha']}")
    else:
        data = json.loads(deck_json.read_text())

    if require_current_sha and data.get("git_sha") != sha:
        raise RuntimeError(f"stale evidence: deck_data git_sha={data.get('git_sha')} != current {sha}")

    quality_path = base / "data" / "quality_report.json"
    tier = 3
    if quality_path.exists():
        tier = int(json.loads(quality_path.read_text()).get("tier", 3))

    data_mode = data["data_mode"]
    watermark = (
        f"{TIER_LABELS.get(tier, f'TIER {tier}')} \u00b7 git {data['git_sha']} "
        f"\u00b7 evidence-generated \u00b7 not investment advice"
    )

    # every screenshot must come from the CURRENT sha evidence dir
    screenshots = sorted(p for pat in ("*.png", "*.jpg") for p in (base / "demo").glob(pat))

    equity: list[float] = []
    eq_csv = base / "historical" / "equity_curve.csv"
    if eq_csv.exists():
        for line in eq_csv.read_text().splitlines()[1:]:
            try:
                equity.append(float(line.split(",")[1]))
            except (IndexError, ValueError):
                continue

    return Evidence(
        data=data,
        base_dir=base,
        sha=data["git_sha"],
        tier=tier,
        data_mode=data_mode,
        watermark=watermark,
        screenshots=screenshots,
        equity_curve=equity,
    )


# ---------------------------------------------------------------------------
# Formatting helpers
# ---------------------------------------------------------------------------


def _pct(x: float | None) -> str:
    return "n/a" if x is None else f"{100 * x:,.2f}%"


def _f2(x: float | None) -> str:
    return "n/a" if x is None else f"{x:,.2f}"


def _render_equity_chart(ev: Evidence) -> Path | None:
    """Render the equity-curve chart from evidence CSV, watermarked on-image."""
    if not ev.equity_curve:
        return None
    try:
        import matplotlib

        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except ImportError:
        return None
    out = ev.base_dir / "pitch" / "equity_curve.png"
    fig, ax = plt.subplots(figsize=(8, 3.6), dpi=120)
    ax.plot(ev.equity_curve, color="#2bd47f", linewidth=1.2)
    ax.set_title(f"Equity curve \u2014 {ev.data_mode} run \u00b7 {len(ev.equity_curve)} steps")
    ax.set_xlabel("step")
    ax.set_ylabel("equity")
    ax.grid(alpha=0.25)
    fig.text(
        0.5,
        0.5,
        ev.watermark,
        ha="center",
        va="center",
        fontsize=11,
        color="gray",
        alpha=0.30,
        rotation=12,
        wrap=True,
    )
    fig.tight_layout()
    fig.savefig(out)
    plt.close(fig)
    return out


# ---------------------------------------------------------------------------
# Slide model: exactly 10 slides, every number from deck_data.json
# ---------------------------------------------------------------------------


def build_slides(ev: Evidence) -> list[dict[str, Any]]:
    d = ev.data
    h = d["historical"]
    s = d["synthetic"]
    minz = d.get("minimized") or {}
    adj = d.get("adjacent_pass") or {}

    is_historical = ev.data_mode in HISTORICAL_MODES
    hist_label = (
        f"Real historical backtest ({ev.data_mode})"
        if is_historical
        else "Synthetic-fixture backtest (NOT historical data)"
    )

    confirmed = sorted(s.get("failed_mechanisms", []))
    failed_str = ", ".join(confirmed) if confirmed else "none confirmed"

    min_metrics = minz.get("metrics") or {}
    adj_metrics = adj.get("metrics") or {}
    adj_found = bool(adj.get("passes"))

    baseline_vs = [
        f"Baseline ({hist_label}): Sharpe {_f2(h['sharpe'])}, max drawdown {_pct(h['max_drawdown'])}, "
        f"cumulative return {_pct(h['cumulative_return'])}.",
        f"Sealed stress search: {s['confirmed_count']} confirmed failures out of {s['evaluated']} worlds "
        f"({_pct(s['failure_rate'])} failure rate).",
        f"Confirmed-failure mechanisms: {failed_str}.",
    ]
    if min_metrics:
        baseline_vs.append(
            f"Worst confirmed stress case ({minz.get('mechanism')}): Sharpe {_f2(min_metrics['sharpe'])}, "
            f"max drawdown {_pct(min_metrics['max_drawdown'])}, cost {_pct(min_metrics['cost_pct_of_capital'])} "
            f"of capital \u2014 invisible to the plain backtest."
        )

    min_bullets = [
        f"Mechanism: {minz.get('mechanism', 'n/a')} \u00b7 seed {minz.get('seed', 'n/a')}.",
        f"Intensity minimized {minz.get('original_intensity', 'n/a')} \u2192 "
        f"{minz.get('minimized_intensity', 'n/a')} and it STILL fails: {minz.get('still_fails', 'n/a')}.",
        f"Violated predicates: {', '.join(minz.get('predicates', [])) or 'n/a'}.",
    ]
    if min_metrics:
        min_bullets.append(
            f"Minimized-case metrics: Sharpe {_f2(min_metrics['sharpe'])}, "
            f"cum. return {_pct(min_metrics['cumulative_return'])}, "
            f"max DD {_pct(min_metrics['max_drawdown'])}."
        )
    if adj_found:
        min_bullets.append(
            f"Adjacent PASSING case: {adj.get('mechanism')} seed delta "
            f"{adj.get('delta_from_failure_seed')} passes (Sharpe {_f2(adj_metrics['sharpe'])}) "
            "\u2014 the failure boundary is sharp and reproducible."
        )
    else:
        min_bullets.append(
            f"Adjacent pass: {adj.get('note', 'no adjacent pass found within search radius')} "
            "\u2014 reported honestly, not fabricated."
        )

    slides: list[dict[str, Any]] = [
        {  # 1
            "title": "Fenrix Strategy Validation Lab",
            "subtitle": TAGLINE,
            "bullets": [
                f"Product flow: {FLOW}",
                f"Strategy hash (immutable): {d['strategy_hash'][:16]}\u2026",
                f"git SHA: {d['git_sha']} \u00b7 data mode of record: {d['data_mode']} \u00b7 universe: {d['universe_size']} assets",
            ],
            "watermark": False,
        },
        {  # 2
            "title": "The problem",
            "subtitle": "A green backtest is a claim, not a validation",
            "bullets": [
                "Backtests only tell you where a strategy already worked \u2014 they are silent about fragility.",
                "Students and retail quants ship strategies validated on one historical path.",
                "Overfitting, cost blindness and regime fragility hide until real capital finds them.",
                "There is no cheap, reproducible way to ask: under which conditions does this break?",
            ],
            "watermark": False,
        },
        {  # 3
            "title": "Product workflow",
            "subtitle": FLOW,
            "bullets": [
                "Describe: plain-English strategy \u2192 structured clause ledger (every clause reviewed).",
                "Review + Lock: mandatory approve step \u2192 immutable version + canonical SHA-256 hash.",
                "Backtest: the SAME locked hash runs a real multi-asset historical backtest.",
                "Stress: the SAME hash enters a sealed synthetic failure search across mechanisms.",
                "Minimize: confirmed failures are shrunk to a minimal reproducible counterexample.",
                "Export: evidence package with manifest, hashes and claim ledger.",
            ],
            "watermark": False,
        },
        {  # 4
            "title": f"Demo run \u2014 {hist_label}",
            "subtitle": f"{d['universe_size']} assets \u00b7 {h['trades']} trades \u00b7 all figures from deck_data.json",
            "kpis": [
                (_pct(h["cumulative_return"]), "cumulative return"),
                (_f2(h["sharpe"]), "Sharpe"),
                (_pct(h["max_drawdown"]), "max drawdown"),
                (_pct(h["volatility"]), "ann. volatility"),
                (_pct(h["cost_pct_of_capital"]), "cost % of capital"),
                (str(h["trades"]), "trades"),
            ],
            "bullets": [
                f"CAGR {_pct(h['cagr'])} \u00b7 Sortino {_f2(h['sortino'])} \u00b7 Calmar {_f2(h['calmar'])}",
                f"Benchmark CAGR {_pct(h['benchmark_cagr'])} \u00b7 information ratio {_f2(h['information_ratio'])}",
                f"Gross exposure avg {_pct(h['gross_exposure_avg'])} \u00b7 net exposure avg {_pct(h['net_exposure_avg'])}",
                "Costs are explicit bounded heuristics (commission/spread/slippage/borrow) \u2014 not broker calibrations.",
            ],
            "image": "equity_curve.png",
            "watermark": True,
        },
        {  # 5
            "title": "What the normal backtest missed",
            "subtitle": "Baseline vs confirmed stress failure",
            "bullets": baseline_vs,
            "watermark": True,
        },
        {  # 6
            "title": "Minimized counterexample",
            "subtitle": "Smallest reproducible world where the locked strategy still fails",
            "bullets": min_bullets,
            "watermark": True,
        },
        {  # 7
            "title": "Education & research value",
            "subtitle": "A falsification lab, not a profit promise",
            "bullets": [
                "Students see WHY a strategy fails \u2014 mechanism, intensity, seed \u2014 not just an equity curve.",
                "Every claim is bound to a hash and an evidence file; grading and review become reproducible.",
                "Instructors can replay the minimized counterexample deterministically.",
                f"Mechanism library searched this run: {', '.join(s['mechanisms_evaluated'])}.",
            ],
            "watermark": False,
        },
        {  # 8
            "title": "Architecture & moat",
            "subtitle": "Immutable contracts + sealed synthetic worlds",
            "bullets": [
                "Clause ledger \u2192 canonical hash: the same immutable artifact flows through every stage.",
                "Sealed synthetic world generator: seeds + mechanisms are reproducible and tamper-evident.",
                "Evidence packages carry SHA-256 manifests binding data, backtest, campaign and replay.",
                "Failure minimizer + adjacent-pass search turn red flags into actionable counterexamples.",
                "Moat: the corpus of confirmed, minimized failures compounds with every run.",
            ],
            "watermark": False,
        },
        {  # 9
            "title": "Honest boundary & roadmap",
            "subtitle": "What this is NOT \u2014 and what comes next",
            "bullets": [
                *d["limitations"],
                "Roadmap: point-in-time fundamentals, broker-calibrated costs, richer mechanism library, live paper-trading bridge.",
            ],
            "watermark": False,
        },
        {  # 10
            "title": "Ask",
            "subtitle": "Help us make strategy falsification the default",
            "bullets": [
                "Pilot with quant-finance classrooms and student funds this semester.",
                "Feedback on mechanism realism from practitioners.",
                "Compute credits for larger sealed stress campaigns.",
                "Reproduce everything: make submission-demo \u00b7 make verify-submission \u00b7 make pitch-deck.",
            ],
            "watermark": False,
        },
    ]
    assert len(slides) == 10, f"deck must have exactly 10 slides, got {len(slides)}"
    return slides


# ---------------------------------------------------------------------------
# HTML renderer
# ---------------------------------------------------------------------------


def render_html(slides: list[dict[str, Any]], ev: Evidence) -> str:
    parts: list[str] = []
    for i, sl in enumerate(slides, 1):
        kpi_html = ""
        if sl.get("kpis"):
            cards = "".join(
                f'<div class="card"><b>{v}</b><span>{label}</span></div>' for v, label in sl["kpis"]
            )
            kpi_html = f'<div class="kpi">{cards}</div>'
        bullets = "".join(f"<li>{b}</li>" for b in sl.get("bullets", []))
        img_html = ""
        if sl.get("image"):
            img_path = ev.base_dir / "pitch" / sl["image"]
            if img_path.exists():
                rel = f"../../../artifacts/submission/{ev.sha}/pitch/{sl['image']}"
                img_html = f'<img src="{rel}" alt="equity curve" style="width:100%;border-radius:8px;margin-top:10px">'
        for shot in ev.screenshots if sl.get("image") else []:
            img_html += (
                f'<img src="../../../{shot}" alt="screenshot" '
                'style="width:100%;border-radius:8px;margin-top:10px">'
            )
        wm = f'<div class="wm">{ev.watermark}</div>' if sl.get("watermark") else ""
        parts.append(
            f'<div class="slide" data-slide="{i}">{wm}<h1>{sl["title"]}</h1>'
            f"<h2>{sl['subtitle']}</h2>{kpi_html}<ul>{bullets}</ul>{img_html}</div>"
        )

    body = "\n".join(parts)
    return f"""<!doctype html>
<html lang="en"><head><meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>Fenrix Strategy Validation Lab \u2014 Submission Deck</title>
<style>
 body{{font-family:-apple-system,Segoe UI,Roboto,Helvetica,Arial,sans-serif;margin:0;background:#0b1020;color:#e7ecf3}}
 .slide{{position:relative;max-width:900px;margin:24px auto;background:#121a2e;border-radius:14px;padding:32px;box-shadow:0 8px 30px rgba(0,0,0,.35)}}
 h1{{font-size:30px;margin:0 0 6px}} h2{{font-size:18px;color:#7fd1ff;margin-top:0;font-weight:500}}
 .kpi{{display:flex;flex-wrap:wrap;gap:14px;margin:14px 0}}
 .card{{flex:1 1 160px;background:#0f1730;border:1px solid #243352;border-radius:10px;padding:14px}}
 .card b{{display:block;font-size:24px;color:#9affc4}} .card span{{font-size:12px;color:#9fb3d1}}
 .wm{{position:absolute;top:10px;right:14px;font-size:11px;color:#ffb020;border:1px solid #ffb020;
      border-radius:6px;padding:3px 8px;opacity:.85;max-width:46%;text-align:right}}
 ul{{line-height:1.55}} code{{background:#0b1020;padding:1px 5px;border-radius:4px;color:#9affc4}}
</style></head><body>
{body}
<p style="text-align:center;color:#5c6b8a;font-size:12px">All numbers rendered from
artifacts/submission/{ev.sha}/pitch/deck_data.json \u00b7 no hand-entered figures.</p>
</body></html>"""


# ---------------------------------------------------------------------------
# PPTX renderer (python-pptx)
# ---------------------------------------------------------------------------


def render_pptx(slides: list[dict[str, Any]], ev: Evidence, out_path: Path) -> Path:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    BG = RGBColor(0x12, 0x1A, 0x2E)
    FG = RGBColor(0xE7, 0xEC, 0xF3)
    ACCENT = RGBColor(0x7F, 0xD1, 0xFF)
    GREEN = RGBColor(0x9A, 0xFF, 0xC4)
    AMBER = RGBColor(0xFF, 0xB0, 0x20)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for sl in slides:
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BG

        tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(1.3))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = sl["title"]
        p.font.size = Pt(34)
        p.font.bold = True
        p.font.color.rgb = FG
        p2 = tf.add_paragraph()
        p2.text = sl["subtitle"]
        p2.font.size = Pt(17)
        p2.font.color.rgb = ACCENT

        top = 1.8
        if sl.get("kpis"):
            n = len(sl["kpis"])
            w = 12.1 / n
            for j, (v, label) in enumerate(sl["kpis"]):
                kb = slide.shapes.add_textbox(Inches(0.6 + j * w), Inches(top), Inches(w - 0.15), Inches(1.1))
                ktf = kb.text_frame
                ktf.word_wrap = True
                kp = ktf.paragraphs[0]
                kp.text = v
                kp.font.size = Pt(22)
                kp.font.bold = True
                kp.font.color.rgb = GREEN
                kl = ktf.add_paragraph()
                kl.text = label
                kl.font.size = Pt(11)
                kl.font.color.rgb = FG
            top += 1.3

        bb = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.1), Inches(7.0 - top))
        btf = bb.text_frame
        btf.word_wrap = True
        for k, b in enumerate(sl.get("bullets", [])):
            bp = btf.paragraphs[0] if k == 0 else btf.add_paragraph()
            bp.text = f"\u2022 {b}"
            bp.font.size = Pt(15)
            bp.font.color.rgb = FG

        if sl.get("image"):
            img = ev.base_dir / "pitch" / sl["image"]
            if img.exists():
                slide.shapes.add_picture(str(img), Inches(6.9), Inches(4.1), width=Inches(6.0))
            for shot in ev.screenshots[:1]:
                slide.shapes.add_picture(str(shot), Inches(0.6), Inches(4.6), width=Inches(5.8))

        if sl.get("watermark"):
            wm = slide.shapes.add_textbox(Inches(0.6), Inches(7.02), Inches(12.1), Inches(0.42))
            wtf = wm.text_frame
            wtf.word_wrap = True
            wp = wtf.paragraphs[0]
            wp.text = ev.watermark
            wp.font.size = Pt(11)
            wp.font.color.rgb = AMBER

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


# ---------------------------------------------------------------------------
# PDF renderer (reportlab)
# ---------------------------------------------------------------------------


def render_pdf(slides: list[dict[str, Any]], ev: Evidence, out_path: Path) -> Path:
    from reportlab.lib.colors import HexColor
    from reportlab.lib.pagesizes import landscape, letter
    from reportlab.lib.utils import simpleSplit
    from reportlab.pdfgen import canvas as rl_canvas

    page = landscape(letter)
    W, H = page
    out_path.parent.mkdir(parents=True, exist_ok=True)
    c = rl_canvas.Canvas(str(out_path), pagesize=page)

    BG = HexColor("#121a2e")
    FG = HexColor("#e7ecf3")
    ACCENT = HexColor("#7fd1ff")
    GREEN = HexColor("#9affc4")
    AMBER = HexColor("#ffb020")

    for sl in slides:
        c.setFillColor(BG)
        c.rect(0, 0, W, H, stroke=0, fill=1)

        c.setFillColor(FG)
        c.setFont("Helvetica-Bold", 26)
        c.drawString(40, H - 60, sl["title"])
        c.setFillColor(ACCENT)
        c.setFont("Helvetica", 14)
        for i, ln in enumerate(simpleSplit(sl["subtitle"], "Helvetica", 14, W - 80)):
            c.drawString(40, H - 85 - i * 16, ln)

        y = H - 130
        if sl.get("kpis"):
            n = len(sl["kpis"])
            w = (W - 80) / n
            for j, (v, label) in enumerate(sl["kpis"]):
                x = 40 + j * w
                c.setFillColor(GREEN)
                c.setFont("Helvetica-Bold", 16)
                c.drawString(x, y, v)
                c.setFillColor(FG)
                c.setFont("Helvetica", 9)
                c.drawString(x, y - 14, label)
            y -= 50

        c.setFillColor(FG)
        c.setFont("Helvetica", 12)
        for b in sl.get("bullets", []):
            for ln in simpleSplit("\u2022 " + b, "Helvetica", 12, W - 80):
                c.drawString(40, y, ln)
                y -= 16
            y -= 4

        if sl.get("image"):
            img = ev.base_dir / "pitch" / sl["image"]
            if img.exists():
                try:
                    c.drawImage(str(img), W - 420, 50, width=380, preserveAspectRatio=True, anchor="sw")
                except Exception:
                    pass

        if sl.get("watermark"):
            c.saveState()
            c.setFillColor(AMBER)
            c.setFont("Helvetica-Bold", 10)
            c.drawString(40, 20, ev.watermark)
            c.setFillColorRGB(0.6, 0.6, 0.6, alpha=0.18)
            c.setFont("Helvetica-Bold", 30)
            c.translate(W / 2, H / 2)
            c.rotate(18)
            c.drawCentredString(0, 0, ev.watermark.split(" \u00b7 ")[0])
            c.restoreState()

        c.showPage()

    c.save()
    return out_path


# ---------------------------------------------------------------------------
# Public entry points
# ---------------------------------------------------------------------------


def build_deck_all(require_current_sha: bool = True) -> dict[str, str]:
    """Build HTML + PPTX + PDF decks from current-SHA evidence. Returns paths."""
    ev = load_evidence(require_current_sha=require_current_sha)
    _render_equity_chart(ev)
    slides = build_slides(ev)

    DECK_DIR.mkdir(parents=True, exist_ok=True)
    html = render_html(slides, ev)
    DECK_PATH.write_text(html)

    pptx_path = render_pptx(slides, ev, DECK_DIR / "fenrix_deck.pptx")
    pdf_path = render_pdf(slides, ev, DECK_DIR / "fenrix_deck.pdf")

    return {
        "html": str(DECK_PATH),
        "pptx": str(pptx_path),
        "pdf": str(pdf_path),
        "evidence": str(ev.base_dir / "pitch" / "deck_data.json"),
        "sha": ev.sha,
        "tier": str(ev.tier),
    }


def build_deck() -> str:
    """Backwards-compatible entry: builds all three formats, returns HTML path."""
    return build_deck_all()["html"]
